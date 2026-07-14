import os
import base64
import uuid
from flask import Flask, request, jsonify
from flask_cors import CORS
from supabase import create_client, Client
from dotenv import load_dotenv
import io
import PyPDF2
import docx
import openpyxl
import pptx
import csv
from datetime import datetime
import pytesseract
from PIL import Image
import requests
from gtts import gTTS
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
# Load environment variables
load_dotenv()

# Get the path to the frontend folder (one level up from backend)
frontend_folder = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))

app = Flask(__name__, static_folder=frontend_folder, static_url_path='')
# Enable CORS for all routes (to allow frontend on port 5500 or similar)
CORS(app)

# Initialize Supabase client
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")

if not SUPABASE_URL or not SUPABASE_KEY:
    raise ValueError("Missing SUPABASE_URL or SUPABASE_KEY in environment variables")

supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)

# ─── REAL TTS ENDPOINT ──────────────────────────────────────────────────────
@app.route('/api/tts', methods=['POST'])
def generate_tts():
    """
    Receives { text: "...", lang: "hi" } and returns a real MP3 audio stream.
    Uses Google's gTTS library - generates a genuine audio file, not a dummy.
    """
    try:
        data = request.json
        text = data.get('text', '').strip()
        lang = data.get('lang', 'en')

        if not text:
            return jsonify({'error': 'No text provided'}), 400

        # Generate real MP3 using gTTS
        tts = gTTS(text=text, lang=lang, slow=False)
        mp3_buffer = io.BytesIO()
        tts.write_to_fp(mp3_buffer)
        mp3_buffer.seek(0)

        from flask import send_file
        return send_file(
            mp3_buffer,
            mimetype='audio/mpeg',
            as_attachment=False,
            download_name='response.mp3'
        )
    except Exception as e:
        print(f'[TTS ERROR] {e}')
        return jsonify({'error': 'NeuroCLI AI Voice Server Error. Please try again.'}), 500
# ─────────────────────────────────────────────────────────────────────────────

# ─── CALL SUMMARY EMAIL ENDPOINT ───────────────────────────────────────────────
@app.route('/api/send-call-summary', methods=['POST'])
def send_call_summary():
    """
    Sends a formatted HTML call summary email to the user's registered email.
    Expects: { to_email, persona_name, conversation, language, duration }
    """
    try:
        data = request.json
        to_email    = data.get('to_email', '').strip()
        persona     = data.get('persona_name', 'AI Agent')
        conversation = data.get('conversation', [])  # [{role, text}]
        language    = data.get('language', 'English')
        duration    = data.get('duration', 'N/A')
        timestamp   = datetime.now().strftime('%B %d, %Y at %I:%M %p')

        if not to_email:
            return jsonify({'error': 'No recipient email'}), 400

        gmail_user = os.getenv('GMAIL_USER', '')
        gmail_pass = os.getenv('GMAIL_APP_PASSWORD', '')

        if not gmail_user or 'your_gmail' in gmail_user:
            return jsonify({'error': 'Email not configured on server'}), 503

        # Build conversation HTML
        convo_html = ''
        for msg in conversation:
            role = msg.get('role', 'user')
            text = msg.get('text', '')
            if role == 'user':
                convo_html += f'''
                <div style="text-align:right;margin:10px 0;">
                    <span style="background:#4f46e5;color:white;padding:8px 14px;border-radius:14px 14px 4px 14px;display:inline-block;max-width:75%;font-size:14px;">{text}</span>
                    <div style="font-size:11px;color:#94a3b8;margin-top:4px;">You</div>
                </div>'''
            else:
                convo_html += f'''
                <div style="text-align:left;margin:10px 0;">
                    <span style="background:#1e2540;color:#e2e8f0;padding:8px 14px;border-radius:14px 14px 14px 4px;display:inline-block;max-width:75%;font-size:14px;border:1px solid rgba(255,255,255,0.06);">{text}</span>
                    <div style="font-size:11px;color:#94a3b8;margin-top:4px;">🤖 {persona}</div>
                </div>'''

        html_body = f"""
        <!DOCTYPE html>
        <html>
        <head><meta charset="UTF-8"></head>
        <body style="margin:0;padding:0;background:#0f172a;font-family:'Segoe UI',Arial,sans-serif;">
            <table width="100%" cellpadding="0" cellspacing="0" style="background:#0f172a;padding:40px 20px;">
                <tr><td align="center">
                    <table width="600" cellpadding="0" cellspacing="0" style="background:#1a2035;border-radius:16px;overflow:hidden;border:1px solid rgba(124,58,237,0.2);">

                        <!-- Header -->
                        <tr><td style="background:linear-gradient(135deg,#7c3aed,#4f46e5);padding:28px 32px;">
                            <h1 style="color:white;margin:0;font-size:22px;">📞 Voice Call Summary</h1>
                            <p style="color:rgba(255,255,255,0.8);margin:6px 0 0;font-size:14px;">{timestamp}</p>
                        </td></tr>

                        <!-- Meta Info -->
                        <tr><td style="padding:24px 32px;border-bottom:1px solid rgba(255,255,255,0.06);">
                            <table width="100%">
                                <tr>
                                    <td style="color:#94a3b8;font-size:13px;padding:4px 0;">🤖 Agent</td>
                                    <td style="color:#e2e8f0;font-size:13px;font-weight:600;text-align:right;">{persona}</td>
                                </tr>
                                <tr>
                                    <td style="color:#94a3b8;font-size:13px;padding:4px 0;">🌐 Language</td>
                                    <td style="color:#e2e8f0;font-size:13px;text-align:right;">{language}</td>
                                </tr>
                                <tr>
                                    <td style="color:#94a3b8;font-size:13px;padding:4px 0;">⏱️ Duration</td>
                                    <td style="color:#e2e8f0;font-size:13px;text-align:right;">{duration}</td>
                                </tr>
                                <tr>
                                    <td style="color:#94a3b8;font-size:13px;padding:4px 0;">💬 Messages</td>
                                    <td style="color:#e2e8f0;font-size:13px;text-align:right;">{len(conversation)}</td>
                                </tr>
                            </table>
                        </td></tr>

                        <!-- Conversation Transcript -->
                        <tr><td style="padding:24px 32px;">
                            <h2 style="color:#a78bfa;font-size:16px;margin:0 0 16px;">📝 Conversation Transcript</h2>
                            <div style="background:#0f172a;border-radius:12px;padding:16px;">
                                {convo_html}
                            </div>
                        </td></tr>

                        <!-- Footer -->
                        <tr><td style="padding:20px 32px;background:#0f172a;border-top:1px solid rgba(255,255,255,0.06);">
                            <p style="color:#64748b;font-size:12px;margin:0;text-align:center;">
                                This summary was automatically generated by <strong style="color:#a78bfa;">NeuroCLI AI Voice Studio</strong>.<br>
                                You are receiving this because you completed a voice call session.
                            </p>
                        </td></tr>

                    </table>
                </td></tr>
            </table>
        </body></html>
        """

        # Send via Gmail SMTP
        msg = MIMEMultipart('alternative')
        msg['Subject'] = f'📞 Your NeuroCLI Call Summary — {persona}'
        msg['From']    = f'NeuroCLI AI <{gmail_user}>'
        msg['To']      = to_email
        msg.attach(MIMEText(html_body, 'html'))

        with smtplib.SMTP_SSL('smtp.gmail.com', 465) as server:
            server.login(gmail_user, gmail_pass)
            server.sendmail(gmail_user, to_email, msg.as_string())

        print(f'[EMAIL] Call summary sent to {to_email}')
        return jsonify({'success': True, 'message': f'Summary sent to {to_email}'})

    except smtplib.SMTPAuthenticationError:
        print('[EMAIL ERROR] Gmail authentication failed. Check GMAIL_USER and GMAIL_APP_PASSWORD.')
        return jsonify({'error': 'Email server authentication failed.'}), 500
    except Exception as e:
        print(f'[EMAIL ERROR] {e}')
        return jsonify({'error': 'Failed to send email summary.'}), 500
# ─────────────────────────────────────────────────────────────────────────────

@app.route('/', methods=['GET'])
def home():
    return app.send_static_file('index.html')

@app.route('/app', methods=['GET'])
@app.route('/chat.html', methods=['GET'])
def chat_app():
    return app.send_static_file('chat.html')

@app.route('/auth', methods=['GET'])
@app.route('/auth.html', methods=['GET'])
def auth_page():
    return app.send_static_file('auth.html')

@app.route('/api/users/sync', methods=['POST'])
def sync_user():
    """
    Syncs a user from Firebase Auth to Supabase Postgres.
    Expects JSON: { "uid": "...", "email": "...", "full_name": "..." }
    """
    try:
        data = request.json
        uid = data.get('uid')
        email = data.get('email')
        full_name = data.get('full_name', '')

        if not uid or not email:
            return jsonify({"error": "Missing uid or email"}), 400

        # Check if user already exists
        response = supabase.table('users').select('*').eq('id', uid).execute()
        
        current_time = datetime.utcnow().isoformat()
        
        if len(response.data) == 0:
            # Insert new user
            supabase.table('users').insert({
                "id": uid,
                "email": email,
                "full_name": full_name,
                "last_login": current_time
            }).execute()
        else:
            # Update existing user (e.g., last_login)
            supabase.table('users').update({
                "full_name": full_name,
                "last_login": current_time
            }).eq('id', uid).execute()

        return jsonify({"success": True, "message": "User synced"}), 200

    except Exception as e:
        print(f"Error syncing user: {e}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/chats/save', methods=['POST'])
def save_chat():
    """
    Saves a chat and its messages.
    Expects JSON:
    {
        "user_id": "...",
        "chat_id": "...",
        "title": "...",
        "messages": [
            { "role": "user", "content": "...", "image_url": "..." },
            ...
        ]
    }
    """
    try:
        data = request.json
        user_id = data.get('user_id')
        chat_id = data.get('chat_id')
        title = data.get('title')
        messages = data.get('messages', [])

        if not user_id or not chat_id:
            return jsonify({"error": "Missing user_id or chat_id"}), 400

        # Upsert the chat metadata
        chat_data = {
            "id": chat_id,
            "user_id": user_id,
            "title": title
        }
        supabase.table('chats').upsert(chat_data).execute()

        # Simple approach: clear old messages for this chat and insert new ones
        # (For production, you'd only insert new ones to save DB writes)
        supabase.table('messages').delete().eq('chat_id', chat_id).execute()

        if messages:
            msg_inserts = []
            for msg in messages:
                msg_inserts.append({
                    "chat_id": chat_id,
                    "role": msg.get('role'),
                    "content": msg.get('content', ''),
                    "image_url": msg.get('image_url')
                })
            supabase.table('messages').insert(msg_inserts).execute()

        return jsonify({"success": True, "message": "Chat saved"}), 200

    except Exception as e:
        print(f"Error saving chat: {e}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/extract-text', methods=['POST'])
def extract_text():
    """
    Extracts text from uploaded documents (PDF, DOCX, XLSX, PPTX, CSV, TXT)
    Expects JSON: { "base64_file": "...", "filename": "example.pdf" }
    """
    try:
        data = request.json
        b64_string = data.get('base64_file')
        filename = data.get('filename', '').lower()
        
        if not b64_string:
            return jsonify({"error": "No base64_file provided"}), 400

        if ',' in b64_string:
            b64_string = b64_string.split(',', 1)[1]
            
        file_data = base64.b64decode(b64_string)
        file_stream = io.BytesIO(file_data)
        extracted_text = ""
        
        try:
            if filename.endswith('.pdf'):
                reader = PyPDF2.PdfReader(file_stream)
                for page in reader.pages:
                    text = page.extract_text()
                    if text: extracted_text += text + "\n"
                    
            elif filename.endswith('.docx'):
                doc = docx.Document(file_stream)
                extracted_text = "\n".join([p.text for p in doc.paragraphs])
                
            elif filename.endswith('.xlsx'):
                wb = openpyxl.load_workbook(file_stream, data_only=True)
                for sheet in wb.worksheets:
                    extracted_text += f"--- Sheet: {sheet.title} ---\n"
                    for row in sheet.iter_rows(values_only=True):
                        extracted_text += ", ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
                        
            elif filename.endswith('.pptx'):
                prs = pptx.Presentation(file_stream)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
                            
            elif filename.endswith('.csv'):
                decoded_str = file_data.decode('utf-8')
                extracted_text = decoded_str
                
            elif filename.endswith('.txt'):
                extracted_text = file_data.decode('utf-8')
                
            else:
                return jsonify({"error": "Unsupported file format"}), 400

            # Limit text length to prevent breaking URL limits
            MAX_CHARS = 10000
            if len(extracted_text) > MAX_CHARS:
                extracted_text = extracted_text[:MAX_CHARS] + "\n...[TRUNCATED]"

            return jsonify({"success": True, "text": extracted_text.strip()}), 200

        except Exception as parse_err:
            print(f"Parse error: {parse_err}")
            return jsonify({"error": "Failed to parse file"}), 500

    except Exception as e:
        print(f"Error extracting text: {e}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/upload', methods=['POST'])
def upload_file():
    """
    Uploads a base64 image string to Supabase Storage 'files' bucket.
    Expects JSON: { "base64_image": "data:image/png;base64,iVBORw0KGgo...", "filename": "optional.png" }
    Returns: { "url": "https://.../files/your-file.png" }
    """
    try:
        data = request.json
        b64_string = data.get('base64_image')
        
        if not b64_string:
            return jsonify({"error": "No base64_image provided"}), 400

        # Parse base64 string
        if ',' in b64_string:
            header, encoded = b64_string.split(',', 1)
            # Extact mime type from header (e.g. data:image/png;base64)
            mime_type = header.split(';')[0].split(':')[1]
            ext = mime_type.split('/')[-1]
        else:
            encoded = b64_string
            mime_type = "application/octet-stream"
            ext = "bin"

        file_data = base64.b64decode(encoded)
        
        # Generate unique filename
        filename = data.get('filename') or f"upload_{uuid.uuid4().hex}.{ext}"
        
        # Upload to Supabase Storage
        # Note: Bucket 'files' must exist and be public
        res = supabase.storage.from_('files').upload(
            path=filename,
            file=file_data,
            file_options={"content-type": mime_type, "x-upsert": "true"}
        )

        # Get public URL
        public_url = supabase.storage.from_('files').get_public_url(filename)
        
        return jsonify({"success": True, "url": public_url}), 200

    except Exception as e:
        print(f"Error uploading file: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/users/analytics', methods=['GET'])
def user_analytics():
    user_id = request.args.get('user_id')
    if not user_id:
        return jsonify({'error': 'Missing user_id'}), 400

    try:
        # Get all chat IDs for this user
        chats_response = supabase.table('chats').select('id').eq('user_id', user_id).execute()
        chat_ids = [chat['id'] for chat in chats_response.data]

        if not chat_ids:
            return jsonify({'success': True, 'metrics': {'day': 0, 'week': 0, 'month': 0}}), 200

        # Fetch messages for these chats
        messages_response = supabase.table('messages') \
            .select('created_at') \
            .in_('chat_id', chat_ids) \
            .eq('role', 'user') \
            .execute()

        from datetime import datetime, timezone
        now = datetime.now(timezone.utc)
        day_count = 0
        week_count = 0
        month_count = 0

        for msg in messages_response.data:
            try:
                # Format: 2026-07-13T04:04:35.163234 or similar
                dt_str = msg['created_at']
                # Try to parse with standard fromisoformat
                dt = datetime.fromisoformat(dt_str.replace('Z', '+00:00'))
                
                # Make naive datetime aware if needed
                if dt.tzinfo is None:
                    dt = dt.replace(tzinfo=timezone.utc)
                    
                delta = now - dt
                if delta.days < 1:
                    day_count += 1
                if delta.days < 7:
                    week_count += 1
                if delta.days < 30:
                    month_count += 1
            except Exception as e:
                pass

        return jsonify({
            'success': True,
            'metrics': {
                'day': day_count,
                'week': week_count,
                'month': month_count
            }
        }), 200

    except Exception as e:
        print("Analytics error:", e)
        return jsonify({'error': str(e)}), 500


@app.route('/api/analyze-scribble', methods=['POST'])
def analyze_scribble():
    try:
        data = request.json
        base64_image = data.get('image')
        
        if not base64_image:
            return jsonify({'error': 'No image provided'}), 400

        api_key = os.getenv('GEMINI_API_KEY')
        if not api_key:
            return jsonify({'error': 'NeuroCLI AI Server Error: Vision Services are currently offline. Please contact support.'}), 500

        prompt = """You are an advanced interactive AI assistant connected to a visual scratchpad. 
The user has hand-drawn something, written a math equation, or sketched a UI layout.
Carefully analyze the image and execute the intent. If it looks like code, write the code. If it looks like math, solve it. If it looks like a UI wireframe description, provide HTML/CSS. If it's just a regular drawing, describe it.

Format your output in beautifully styled Markdown."""

        payload = {
            "contents": [{
                "parts": [
                    { "text": prompt },
                    {
                        "inlineData": {
                            "mimeType": "image/jpeg",
                            "data": base64_image
                        }
                    }
                ]
            }]
        }

        response = requests.post(
            f'https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={api_key}',
            headers={'Content-Type': 'application/json'},
            json=payload
        )

        if not response.ok:
            error_data = response.json()
            # Log the actual error for the developer to see in the terminal
            print(f"Backend API Error: {error_data}")
            return jsonify({'error': 'NeuroCLI AI Server Error: Unable to process vision request at this time. Please try again later.'}), response.status_code

        response_data = response.json()
        ai_text = response_data['candidates'][0]['content']['parts'][0]['text']

        return jsonify({'markdown': ai_text})

    except Exception as e:
        print(f"Scribble analysis error: {e}")
        return jsonify({'error': str(e)}), 500

# ─── SAVE TO DISK ENDPOINT (Hologram Editor) ──────────────────────────────────
@app.route('/api/save_file', methods=['POST'])
def save_file_to_disk():
    """
    Saves the provided code/text directly to the user's local disk.
    Expects: { filename: "script.py", content: "print('Hello')" }
    """
    try:
        data = request.json
        filename = data.get('filename', 'untitled.txt')
        content = data.get('content', '')

        if not filename or not content:
            return jsonify({'error': 'Filename and content are required'}), 400

        # Save to the root project directory for demonstration
        save_path = os.path.join(frontend_folder, filename)
        
        with open(save_path, 'w', encoding='utf-8') as f:
            f.write(content)

        return jsonify({'success': True, 'path': save_path, 'message': f'File saved successfully to {save_path}'}), 200

    except Exception as e:
        print(f'[SAVE ERROR] {e}')
        return jsonify({'error': str(e)}), 500
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == '__main__':
    print("Starting NeuroCLI AI Backend on port 5007...")
    app.run(host='0.0.0.0', port=5007, debug=True)

